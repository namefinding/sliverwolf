from __future__ import annotations

import fnmatch
from difflib import SequenceMatcher
import os
import re
import shutil
import subprocess
import xml.etree.ElementTree as ET
from enum import Enum
from mimetypes import guess_type
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pydantic import AliasChoices, BaseModel, Field, model_validator

from local_agent.protocol.models import OutputKind, ToolManifest
from local_agent.utils.file_query_normalizer import FileQueryNormalizer
from local_agent.utils.workspace_path import WorkspacePathNormalizer


class ListFilesInput(BaseModel):
    path: str
    recursive: bool = False
    include_dirs: bool = True
    patterns: list[str] = Field(default_factory=list)


class ReadFilesInput(BaseModel):
    paths: list[str]
    encoding: str = "utf-8"
    max_bytes: int = 200_000


class MatchMode(str, Enum):
    ANY = "any"
    ALL = "all"


class SearchTextInput(BaseModel):
    path: str
    query: str | None = None
    terms: list[str] = Field(default_factory=list)
    match_mode: MatchMode = MatchMode.ANY
    recursive: bool = True
    patterns: list[str] = Field(default_factory=lambda: ["*.py", "*.md", "*.txt", "*.json", "*.yaml", "*.yml"])
    max_matches: int = 20

    @model_validator(mode="after")
    def validate_search_source(self) -> "SearchTextInput":
        if not self.query and not self.terms:
            raise ValueError("query or terms must be provided")
        if self.query and not self.terms:
            self.terms = [term for term in self.query.split() if term.strip()]
        return self


class WriteFileInput(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"
    overwrite: bool = True


class WriteDocxInput(BaseModel):
    path: str
    title: str | None = None
    content: str = ""
    paragraphs: list[str] = Field(default_factory=list)
    overwrite: bool = True


class WriteXlsxInput(BaseModel):
    path: str
    title: str | None = None
    content: str = ""
    rows: list[list[str | int | float | bool | None]] = Field(default_factory=list)
    sheet_name: str = "Sheet1"
    overwrite: bool = True


class WritePptxInput(BaseModel):
    path: str
    title: str | None = None
    content: str = ""
    bullets: list[str] = Field(default_factory=list)
    overwrite: bool = True


class SearchByNameInput(BaseModel):
    path: str = "."
    query: str
    query_terms: list[str] = Field(default_factory=list)
    alias_terms: list[str] = Field(default_factory=list)
    recursive: bool = True
    scope_mode: str = "subtree"
    target_kind: str = "any"
    extensions: list[str] = Field(default_factory=list)
    include_dirs: bool = True
    top_k: int = 8


class AppendFileInput(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"
    create: bool = True


class MetadataInput(BaseModel):
    path: str


class PreviewInput(BaseModel):
    path: str
    encoding: str = "utf-8"
    max_chars: int = 1200
    max_children: int = 12


class ExtractTextInput(BaseModel):
    paths: list[str]
    encoding: str = "utf-8"
    max_chars: int = 12000
    max_rows_per_sheet: int = 10


class ExtractStructureInput(BaseModel):
    path: str
    encoding: str = "utf-8"
    include_text: bool = True
    max_blocks: int = 200
    max_chars_per_block: int = 1200
    max_rows_per_sheet: int = 10


class SearchBlocksInput(BaseModel):
    path: str
    query: str | None = None
    terms: list[str] = Field(default_factory=list)
    max_matches: int = 8
    max_blocks: int = 200
    max_chars_per_block: int = 1200

    @model_validator(mode="after")
    def validate_search_source(self) -> "SearchBlocksInput":
        if not self.query and not self.terms:
            raise ValueError("query or terms must be provided")
        if self.query and not self.terms:
            self.terms = [term for term in FileModule._tokenize_search_query(self.query) if term.strip()]
        return self


class DocxEditItem(BaseModel):
    block_id: str
    action: str = "replace"
    text: str = ""


class EditDocxInput(BaseModel):
    source_path: str
    output_path: str
    edits: list[DocxEditItem]
    overwrite: bool = True


class RenderDocxFromTemplateInput(BaseModel):
    template_path: str
    output_path: str
    source_path: str | None = None
    content: str = ""
    paragraphs: list[str] = Field(default_factory=list)
    title: str | None = None
    overwrite: bool = True

    @model_validator(mode="after")
    def validate_render_source(self) -> "RenderDocxFromTemplateInput":
        if self.paragraphs:
            return self
        if self.content.strip():
            return self
        if isinstance(self.source_path, str) and self.source_path.strip():
            return self
        raise ValueError("source_path, content, or paragraphs must be provided")


class MakeDirectoryInput(BaseModel):
    path: str
    exist_ok: bool = True
    parents: bool = True


class CopyPathInput(BaseModel):
    src_path: str
    dest_path: str
    overwrite: bool = False


class MovePathInput(BaseModel):
    src_path: str
    dest_path: str
    overwrite: bool = False


class RenamePathInput(BaseModel):
    path: str
    new_name: str
    overwrite: bool = False


class DeletePathInput(BaseModel):
    path: str
    recursive: bool = False
    missing_ok: bool = False


class OpenPathInput(BaseModel):
    path: str


class PathListInput(BaseModel):
    paths: list[str]
    continue_on_error: bool = False


class WriteManyItem(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"
    overwrite: bool = True


class WriteManyInput(BaseModel):
    items: list[WriteManyItem]
    continue_on_error: bool = False


class AppendManyItem(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"
    create: bool = True


class AppendManyInput(BaseModel):
    items: list[AppendManyItem]
    continue_on_error: bool = False


class PreviewManyInput(BaseModel):
    paths: list[str]
    encoding: str = "utf-8"
    max_chars: int = 1200
    max_children: int = 12
    continue_on_error: bool = False


class MakeDirectoryManyInput(BaseModel):
    paths: list[str]
    exist_ok: bool = True
    parents: bool = True
    continue_on_error: bool = False


class CopyManyItem(BaseModel):
    src_path: str = Field(validation_alias=AliasChoices("src_path", "source_path"))
    dest_path: str
    overwrite: bool = False


class CopyManyInput(BaseModel):
    items: list[CopyManyItem]
    continue_on_error: bool = False


class MoveManyItem(BaseModel):
    src_path: str = Field(validation_alias=AliasChoices("src_path", "source_path"))
    dest_path: str
    overwrite: bool = False


class MoveManyInput(BaseModel):
    items: list[MoveManyItem]
    continue_on_error: bool = False


class RenameManyItem(BaseModel):
    path: str
    new_name: str
    overwrite: bool = False


class RenameManyInput(BaseModel):
    items: list[RenameManyItem]
    continue_on_error: bool = False


class DeleteManyInput(BaseModel):
    paths: list[str]
    recursive: bool = False
    missing_ok: bool = False
    continue_on_error: bool = False


class FileModule:
    def __init__(self, workspace_root: str) -> None:
        self.workspace_root = Path(workspace_root).resolve()
        self.path_normalizer = WorkspacePathNormalizer(str(self.workspace_root))

    def manifests(self) -> list[ToolManifest]:
        return [
            ToolManifest(
                tool_name="file.list",
                module="file",
                description="List files or directories under a target path.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.DIRECTORY_ENTRIES],
                input_schema=ListFilesInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"entries": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.search_by_name",
                module="file",
                description="Search files or folders by file name, title-like terms, and path similarity.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_CANDIDATES],
                input_schema=SearchByNameInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"candidates": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.read",
                module="file",
                description="Read one or more files.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.FILE_CONTENTS],
                input_schema=ReadFilesInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"files": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.extract_text",
                module="file",
                description="Extract readable text from text files and Office documents such as docx, pptx, and xlsx.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.FILE_CONTENTS],
                input_schema=ExtractTextInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"files": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.extract_structure",
                module="file",
                description="Extract structured blocks from a local file, especially docx documents, including paragraphs, headings, and table cells.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=ExtractStructureInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"blocks": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.search_blocks",
                module="file",
                description="Search structured blocks inside a local document and return the most relevant sections.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=SearchBlocksInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"matches": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.search_text",
                module="file",
                description="Search text within files under a path using one or more search terms.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.SEARCH_MATCHES],
                input_schema=SearchTextInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"matches": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.write",
                module="file",
                description="Write a file inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=WriteFileInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.write_docx",
                module="file",
                description="Write a Word document (.docx) inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=WriteDocxInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.edit_docx",
                module="file",
                description="Apply targeted block-level edits to a Word document (.docx) while keeping the surrounding document structure.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=EditDocxInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.render_docx_from_template",
                module="file",
                description="Render a new Word document from a template file and source content while reusing the template's overall layout.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=RenderDocxFromTemplateInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.write_xlsx",
                module="file",
                description="Write an Excel workbook (.xlsx) inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=WriteXlsxInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.write_pptx",
                module="file",
                description="Write a PowerPoint presentation (.pptx) inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=WritePptxInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.write_many",
                module="file",
                description="Write multiple files inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=WriteManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.append",
                module="file",
                description="Append text to a file inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=AppendFileInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.append_many",
                module="file",
                description="Append text to multiple files inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.FILE_WRITTEN],
                input_schema=AppendManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.metadata",
                module="file",
                description="Read metadata for a file or folder inside the workspace root.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=MetadataInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.metadata_many",
                module="file",
                description="Read metadata for multiple files or folders inside the workspace root.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=PathListInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"items": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.preview",
                module="file",
                description="Preview a file or folder without reading the full contents.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=PreviewInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.preview_many",
                module="file",
                description="Preview multiple files or folders without reading the full contents.",
                side_effect=False,
                idempotent=True,
                produces=[OutputKind.OBJECT_DETAILS],
                input_schema=PreviewManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"items": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.mkdir",
                module="file",
                description="Create a directory inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_CREATED],
                input_schema=MakeDirectoryInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.mkdir_many",
                module="file",
                description="Create multiple directories inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_CREATED],
                input_schema=MakeDirectoryManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.copy",
                module="file",
                description="Copy a file or folder to a new location inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_CREATED],
                input_schema=CopyPathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.copy_many",
                module="file",
                description="Copy multiple files or folders to new locations inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_CREATED],
                input_schema=CopyManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.move",
                module="file",
                description="Move a file or folder to a new location inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_UPDATED],
                input_schema=MovePathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.move_many",
                module="file",
                description="Move multiple files or folders to new locations inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_UPDATED],
                input_schema=MoveManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.rename",
                module="file",
                description="Rename a file or folder inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_UPDATED],
                input_schema=RenamePathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.rename_many",
                module="file",
                description="Rename multiple files or folders inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_UPDATED],
                input_schema=RenameManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.delete",
                module="file",
                description="Delete a file or folder inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_DELETED],
                input_schema=DeletePathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.delete_many",
                module="file",
                description="Delete multiple files or folders inside the workspace root.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_DELETED],
                input_schema=DeleteManyInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.open_path",
                module="file",
                description="Open a file or folder with the system default handler.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_OPENED],
                input_schema=OpenPathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.open_many",
                module="file",
                description="Open multiple files or folders with the system default handler.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_OPENED],
                input_schema=PathListInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
            ToolManifest(
                tool_name="file.reveal_in_explorer",
                module="file",
                description="Reveal a file or folder in the system file explorer.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_OPENED],
                input_schema=OpenPathInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"path": {"type": "string"}}},
            ),
            ToolManifest(
                tool_name="file.reveal_many",
                module="file",
                description="Reveal multiple files or folders in the system file explorer.",
                side_effect=True,
                idempotent=False,
                produces=[OutputKind.PATH_OPENED],
                input_schema=PathListInput.model_json_schema(),
                output_schema={"type": "object", "properties": {"paths": {"type": "array"}, "results": {"type": "array"}}},
            ),
        ]

    def executor_map(self) -> dict[str, Any]:
        return {
            "file.list": self.list_files,
            "file.search_by_name": self.search_by_name,
            "file.read": self.read_files,
            "file.extract_text": self.extract_text,
            "file.extract_structure": self.extract_structure,
            "file.search_blocks": self.search_blocks,
            "file.search_text": self.search_text,
            "file.write": self.write_file,
            "file.write_docx": self.write_docx,
            "file.edit_docx": self.edit_docx,
            "file.render_docx_from_template": self.render_docx_from_template,
            "file.write_xlsx": self.write_xlsx,
            "file.write_pptx": self.write_pptx,
            "file.write_many": self.write_many,
            "file.append": self.append_file,
            "file.append_many": self.append_many,
            "file.metadata": self.metadata,
            "file.metadata_many": self.metadata_many,
            "file.preview": self.preview,
            "file.preview_many": self.preview_many,
            "file.mkdir": self.make_directory,
            "file.mkdir_many": self.make_directory_many,
            "file.copy": self.copy_path,
            "file.copy_many": self.copy_many,
            "file.move": self.move_path,
            "file.move_many": self.move_many,
            "file.rename": self.rename_path,
            "file.rename_many": self.rename_many,
            "file.delete": self.delete_path,
            "file.delete_many": self.delete_many,
            "file.open_path": self.open_path,
            "file.open_many": self.open_many,
            "file.reveal_in_explorer": self.reveal_in_explorer,
            "file.reveal_many": self.reveal_many,
        }

    def resolve_path(self, raw_path: str) -> Path:
        return self.path_normalizer.resolve(raw_path)

    def ensure_workspace_path(self, target: Path) -> None:
        try:
            target.relative_to(self.workspace_root)
        except ValueError as exc:
            raise PermissionError(f"Path is outside workspace: {target}") from exc

    def list_files(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = ListFilesInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        iterator = target.rglob("*") if payload.recursive else target.iterdir()

        entries = []
        for item in iterator:
            if not payload.include_dirs and item.is_dir():
                continue
            if payload.patterns and not any(fnmatch.fnmatch(item.name, pattern) for pattern in payload.patterns):
                continue
            entries.append({"path": str(item), "name": item.name, "is_dir": item.is_dir()})
        return {"entries": entries}

    def search_by_name(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = SearchByNameInput.model_validate(arguments)
        root = self.resolve_path(payload.path)
        self.ensure_workspace_path(root)
        iterator = root.rglob("*") if payload.recursive else root.iterdir()
        normalized_query = FileQueryNormalizer.normalize(payload.query)
        query_text = normalized_query.normalized_text or payload.query.strip().lower()
        query_terms = payload.query_terms or normalized_query.core_terms or [term for term in self._tokenize_search_query(payload.query) if term]
        alias_terms = payload.alias_terms or normalized_query.alias_terms
        normalized_exts = {
            ext.lower() if ext.startswith(".") else f".{ext.lower()}"
            for ext in payload.extensions
            if isinstance(ext, str) and ext.strip()
        }
        file_type_hints = normalized_query.file_type_hints or list(normalized_exts)

        candidates: list[dict[str, Any]] = []
        for item in iterator:
            is_dir = item.is_dir()
            if is_dir and not payload.include_dirs:
                continue
            if payload.target_kind == "file" and is_dir:
                continue
            if payload.target_kind == "folder" and not is_dir:
                continue
            if normalized_exts and item.is_file() and item.suffix.lower() not in normalized_exts:
                continue

            score, reason = self._score_name_candidate(
                item,
                query_text=query_text,
                query_terms=query_terms,
                alias_terms=alias_terms,
                file_type_hints=file_type_hints,
            )
            if score <= 0:
                continue
            candidates.append(
                {
                    "path": str(item),
                    "name": item.name,
                    "is_dir": is_dir,
                    "score": round(self._apply_name_scope_bonus(score, item, root, payload.scope_mode), 4),
                    "match_reason": reason,
                }
            )

        candidates.sort(key=lambda item: (item["score"], not item["is_dir"], item["name"].lower()), reverse=True)
        return {"query": payload.query, "candidates": candidates[: payload.top_k]}

    @staticmethod
    def _apply_name_scope_bonus(score: float, item: Path, root: Path, scope_mode: str) -> float:
        if score <= 0 or str(scope_mode or "").strip().lower() != "shallow_first":
            return score
        try:
            relative = item.resolve().relative_to(root.resolve())
            depth = max(len(relative.parts) - 1, 0)
        except Exception:
            depth = 99
        return min(1.0, score + max(0.0, 0.18 - min(depth, 6) * 0.04))

    def read_files(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = ReadFilesInput.model_validate(arguments)
        files = []
        for raw_path in payload.paths:
            path = self.resolve_path(raw_path)
            self.ensure_workspace_path(path)
            content = path.read_text(encoding=payload.encoding)
            trimmed = content[: payload.max_bytes]
            files.append(
                {
                    "path": str(path),
                    "size": path.stat().st_size,
                    "content": trimmed,
                    "truncated": len(content) > len(trimmed),
                }
            )
        return {"files": files}

    def extract_text(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = ExtractTextInput.model_validate(arguments)
        files = []
        for raw_path in payload.paths:
            path = self.resolve_path(raw_path)
            self.ensure_workspace_path(path)
            if not path.exists():
                raise FileNotFoundError(f"Path does not exist: {path}")
            if not path.is_file():
                raise IsADirectoryError(f"Path must be a file: {path}")
            content, extraction_type = self._extract_text_from_file(
                path,
                encoding=payload.encoding,
                max_chars=payload.max_chars,
                max_rows_per_sheet=payload.max_rows_per_sheet,
            )
            files.append(
                {
                    "path": str(path),
                    "size": path.stat().st_size,
                    "content": content[: payload.max_chars],
                    "truncated": len(content) > payload.max_chars,
                    "extraction_type": extraction_type,
                }
            )
        return {"files": files}

    def extract_structure(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = ExtractStructureInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if not path.exists():
            raise FileNotFoundError(f"Path does not exist: {path}")
        if not path.is_file():
            raise IsADirectoryError(f"Path must be a file: {path}")
        structure = self._extract_structure_from_file(
            path,
            encoding=payload.encoding,
            include_text=payload.include_text,
            max_blocks=payload.max_blocks,
            max_chars_per_block=payload.max_chars_per_block,
            max_rows_per_sheet=payload.max_rows_per_sheet,
        )
        structure["path"] = str(path)
        return structure

    def search_blocks(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = SearchBlocksInput.model_validate(arguments)
        structure = self.extract_structure(
            {
                "path": payload.path,
                "include_text": True,
                "max_blocks": payload.max_blocks,
                "max_chars_per_block": payload.max_chars_per_block,
            }
        )
        query_terms = [term.lower() for term in payload.terms if term.strip()]
        query_text = (payload.query or " ".join(payload.terms)).strip().lower()
        matches: list[dict[str, Any]] = []
        for block in structure.get("blocks", []):
            if not isinstance(block, dict):
                continue
            block_text = str(block.get("text", "")).strip()
            if not block_text:
                continue
            lowered_text = block_text.lower()
            style_name = str(block.get("style_name") or "").lower()
            score = 0.0
            if query_text and query_text in lowered_text:
                score += 0.75
            term_hits = sum(1 for term in query_terms if term in lowered_text)
            if term_hits:
                score += 0.2 + 0.5 * (term_hits / max(len(query_terms), 1))
            if term_hits and "heading" in style_name:
                score += 0.08
            if score <= 0:
                continue
            matches.append(
                {
                    **block,
                    "score": round(min(score, 1.0), 4),
                    "match_reason": "query_match" if query_text and query_text in lowered_text else "term_overlap",
                }
            )
        matches.sort(key=lambda item: item.get("score", 0.0), reverse=True)
        return {
            "path": structure.get("path", payload.path),
            "file_type": structure.get("file_type"),
            "query": payload.query,
            "terms": payload.terms,
            "matches": matches[: payload.max_matches],
        }

    def search_text(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = SearchTextInput.model_validate(arguments)
        root = self.resolve_path(payload.path)
        self.ensure_workspace_path(root)
        iterator = root.rglob("*") if payload.recursive else root.iterdir()
        search_terms = [term.lower() for term in payload.terms if term.strip()]

        matches = []
        for item in iterator:
            if not item.is_file():
                continue
            if payload.patterns and not any(fnmatch.fnmatch(item.name, pattern) for pattern in payload.patterns):
                continue
            try:
                text = item.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                continue
            for line_no, line in enumerate(text.splitlines(), start=1):
                lowered_line = line.lower()
                if payload.match_mode == MatchMode.ALL:
                    is_match = all(term in lowered_line for term in search_terms)
                else:
                    is_match = any(term in lowered_line for term in search_terms)
                if is_match:
                    matches.append({"path": str(item), "line": line_no, "snippet": line.strip()})
                    if len(matches) >= payload.max_matches:
                        return {"matches": matches}
        return {"matches": matches}

    def write_file(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = WriteFileInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if path.exists() and not payload.overwrite:
            raise FileExistsError(f"File already exists: {path}")
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(payload.content, encoding=payload.encoding)
        return {"path": str(path), "bytes_written": len(payload.content.encode(payload.encoding))}

    def write_docx(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = WriteDocxInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if path.exists() and not payload.overwrite:
            raise FileExistsError(f"File already exists: {path}")
        path.parent.mkdir(parents=True, exist_ok=True)

        from docx import Document

        document = Document()
        title = (payload.title or path.stem).strip()
        if title:
            document.add_heading(title, level=1)
        paragraphs = payload.paragraphs or self._split_paragraphs(payload.content)
        if not paragraphs and payload.content.strip():
            paragraphs = [payload.content.strip()]
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        document.save(path)
        return {"path": str(path), "bytes_written": path.stat().st_size}

    def edit_docx(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = EditDocxInput.model_validate(arguments)
        source_path = self.resolve_path(payload.source_path)
        output_path = self.resolve_path(payload.output_path)
        self.ensure_workspace_path(source_path)
        self.ensure_workspace_path(output_path)
        if not source_path.exists():
            raise FileNotFoundError(f"Source path does not exist: {source_path}")
        if source_path.suffix.lower() != ".docx" or output_path.suffix.lower() != ".docx":
            raise ValueError("file.edit_docx requires .docx source_path and output_path")
        if output_path.exists() and not payload.overwrite and output_path != source_path:
            raise FileExistsError(f"File already exists: {output_path}")

        from docx import Document

        document = Document(source_path)
        block_refs = self._build_docx_block_reference_map(document)
        applied: list[dict[str, Any]] = []
        for edit in payload.edits:
            block_ref = block_refs.get(edit.block_id)
            if block_ref is None:
                raise ValueError(f"Unknown block_id: {edit.block_id}")
            action = edit.action.strip().lower()
            if action == "replace":
                self._replace_docx_block_text(block_ref, edit.text)
            elif action == "delete":
                self._delete_docx_block(block_ref)
            elif action == "insert_after":
                self._insert_docx_block_after(block_ref, edit.text)
            else:
                raise ValueError(f"Unsupported edit action: {edit.action}")
            applied_item = {"block_id": edit.block_id, "action": action}
            preview = self._preview_docx_edit_text(edit.text)
            if preview:
                applied_item["text_preview"] = preview
            applied.append(applied_item)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return {
            "path": str(output_path),
            "source_path": str(source_path),
            "edited": True,
            "edit_count": len(applied),
            "applied_edits": applied,
        }

    @staticmethod
    def _preview_docx_edit_text(text: str, max_chars: int = 140) -> str:
        normalized = " ".join(str(text or "").split())
        if not normalized:
            return ""
        if len(normalized) <= max_chars:
            return normalized
        return normalized[: max_chars - 1].rstrip() + "…"

    def render_docx_from_template(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = RenderDocxFromTemplateInput.model_validate(arguments)
        template_path = self.resolve_path(payload.template_path)
        output_path = self.resolve_path(payload.output_path)
        self.ensure_workspace_path(template_path)
        self.ensure_workspace_path(output_path)
        if not template_path.exists():
            raise FileNotFoundError(f"Template path does not exist: {template_path}")
        if template_path.suffix.lower() != ".docx" or output_path.suffix.lower() != ".docx":
            raise ValueError("file.render_docx_from_template requires .docx template_path and output_path")
        if output_path.exists() and not payload.overwrite:
            raise FileExistsError(f"File already exists: {output_path}")

        source_path: Path | None = None
        if isinstance(payload.source_path, str) and payload.source_path.strip():
            source_path = self.resolve_path(payload.source_path)
            self.ensure_workspace_path(source_path)
            if not source_path.exists():
                raise FileNotFoundError(f"Source path does not exist: {source_path}")

        paragraphs = self._resolve_render_docx_paragraphs(payload, source_path)

        from docx import Document

        document = Document(template_path)
        template_paragraphs = document.paragraphs
        if not template_paragraphs:
            template_paragraphs = [document.add_paragraph("")]

        content_index = 0
        if payload.title:
            template_paragraphs[0].text = payload.title.strip()
            content_index = 1
        elif paragraphs:
            template_paragraphs[0].text = paragraphs[0]
            paragraphs = paragraphs[1:]
            content_index = 1

        reusable_style = template_paragraphs[max(content_index - 1, 0)].style if template_paragraphs else None
        for paragraph in template_paragraphs[content_index:]:
            if paragraphs:
                paragraph.text = paragraphs.pop(0)
            else:
                paragraph.text = ""
        for extra_text in paragraphs:
            new_paragraph = document.add_paragraph(extra_text)
            if reusable_style is not None:
                try:
                    new_paragraph.style = reusable_style
                except Exception:  # noqa: BLE001
                    pass

        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return {
            "path": str(output_path),
            "template_path": str(template_path),
            "source_path": None if source_path is None else str(source_path),
            "rendered": True,
            "paragraph_count": len(document.paragraphs),
        }

    def write_xlsx(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = WriteXlsxInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if path.exists() and not payload.overwrite:
            raise FileExistsError(f"File already exists: {path}")
        path.parent.mkdir(parents=True, exist_ok=True)

        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = (payload.sheet_name or "Sheet1")[:31]
        title = (payload.title or path.stem).strip()
        if title:
            sheet.append([title])
        rows = payload.rows or self._build_tabular_rows(payload.content)
        if not rows and payload.content.strip():
            rows = [["content", payload.content.strip()]]
        for row in rows:
            sheet.append(list(row))
        workbook.save(path)
        return {"path": str(path), "bytes_written": path.stat().st_size}

    def write_pptx(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = WritePptxInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if path.exists() and not payload.overwrite:
            raise FileExistsError(f"File already exists: {path}")
        path.parent.mkdir(parents=True, exist_ok=True)

        from pptx import Presentation

        presentation = Presentation()
        title = (payload.title or path.stem).strip()
        bullets = payload.bullets or self._build_bullets(payload.content)

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title or path.stem
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        if bullets:
            first = text_frame.paragraphs[0]
            first.text = bullets[0]
            for bullet in bullets[1:]:
                para = text_frame.add_paragraph()
                para.text = bullet
                para.level = 0
        elif payload.content.strip():
            text_frame.paragraphs[0].text = payload.content.strip()
        presentation.save(path)
        return {"path": str(path), "bytes_written": path.stat().st_size}

    def write_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = WriteManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.items,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda item: item.path,
            operation=lambda item: self.write_file(item.model_dump()),
        )

    def append_file(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = AppendFileInput.model_validate(arguments)
        path = self.resolve_path(payload.path)
        self.ensure_workspace_path(path)
        if not path.exists() and not payload.create:
            raise FileNotFoundError(f"File does not exist: {path}")
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("a", encoding=payload.encoding) as handle:
            handle.write(payload.content)
        return {"path": str(path), "bytes_written": len(payload.content.encode(payload.encoding)), "appended": True}

    def append_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = AppendManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.items,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda item: item.path,
            operation=lambda item: self.append_file(item.model_dump()),
        )

    def metadata(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = MetadataInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        if not target.exists():
            raise FileNotFoundError(f"Path does not exist: {target}")
        return self._metadata_for_path(target)

    def metadata_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = PathListInput.model_validate(arguments)
        result = self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.metadata({"path": raw_path}),
        )
        result["items"] = result["results"]
        return result

    def preview(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = PreviewInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        if not target.exists():
            raise FileNotFoundError(f"Path does not exist: {target}")
        base = self._metadata_for_path(target)
        if target.is_dir():
            children = []
            for item in sorted(target.iterdir(), key=lambda entry: (not entry.is_dir(), entry.name.lower()))[: payload.max_children]:
                children.append({"name": item.name, "path": str(item), "is_dir": item.is_dir()})
            base["children"] = children
            base["preview_text"] = ""
            return base

        text_preview = ""
        binary = False
        try:
            text_preview = target.read_text(encoding=payload.encoding)[: payload.max_chars]
        except UnicodeDecodeError:
            binary = True
        base["preview_text"] = text_preview
        base["binary"] = binary
        return base

    def preview_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = PreviewManyInput.model_validate(arguments)
        result = self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.preview(
                {
                    "path": raw_path,
                    "encoding": payload.encoding,
                    "max_chars": payload.max_chars,
                    "max_children": payload.max_children,
                }
            ),
        )
        result["items"] = result["results"]
        return result

    def make_directory(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = MakeDirectoryInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        target.mkdir(parents=payload.parents, exist_ok=payload.exist_ok)
        return {"path": str(target), "created": True, "is_dir": True}

    def make_directory_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = MakeDirectoryManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.make_directory(
                {"path": raw_path, "exist_ok": payload.exist_ok, "parents": payload.parents}
            ),
        )

    def copy_path(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = CopyPathInput.model_validate(arguments)
        src = self.resolve_path(payload.src_path)
        dest = self.resolve_path(payload.dest_path)
        self.ensure_workspace_path(src)
        self.ensure_workspace_path(dest)
        if not src.exists():
            raise FileNotFoundError(f"Source path does not exist: {src}")
        if dest.exists():
            if not payload.overwrite:
                raise FileExistsError(f"Destination already exists: {dest}")
            self._remove_existing_path(dest)
        dest.parent.mkdir(parents=True, exist_ok=True)
        if src.is_dir():
            shutil.copytree(src, dest)
        else:
            shutil.copy2(src, dest)
        return {"path": str(dest), "source_path": str(src), "copied": True}

    def copy_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = CopyManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.items,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda item: item.dest_path,
            operation=lambda item: self.copy_path(item.model_dump()),
        )

    def move_path(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = MovePathInput.model_validate(arguments)
        src = self.resolve_path(payload.src_path)
        dest = self.resolve_path(payload.dest_path)
        self.ensure_workspace_path(src)
        self.ensure_workspace_path(dest)
        if not src.exists():
            raise FileNotFoundError(f"Source path does not exist: {src}")
        if src == self.workspace_root:
            raise PermissionError("Cannot move the workspace root")
        if dest.exists():
            if not payload.overwrite:
                raise FileExistsError(f"Destination already exists: {dest}")
            self._remove_existing_path(dest)
        dest.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dest))
        return {"path": str(dest), "source_path": str(src), "moved": True}

    def move_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = MoveManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.items,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda item: item.dest_path,
            operation=lambda item: self.move_path(item.model_dump()),
        )

    def rename_path(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = RenamePathInput.model_validate(arguments)
        src = self.resolve_path(payload.path)
        self.ensure_workspace_path(src)
        if not src.exists():
            raise FileNotFoundError(f"Path does not exist: {src}")
        if src == self.workspace_root:
            raise PermissionError("Cannot rename the workspace root")
        dest = src.with_name(payload.new_name)
        self.ensure_workspace_path(dest)
        if dest.exists():
            if not payload.overwrite:
                raise FileExistsError(f"Destination already exists: {dest}")
            self._remove_existing_path(dest)
        src.rename(dest)
        return {"path": str(dest), "source_path": str(src), "renamed": True}

    def rename_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = RenameManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.items,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda item: item.new_name,
            operation=lambda item: self.rename_path(item.model_dump()),
        )

    def delete_path(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = DeletePathInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        if target == self.workspace_root:
            raise PermissionError("Cannot delete the workspace root")
        if not target.exists():
            if payload.missing_ok:
                return {"path": str(target), "deleted": False, "missing": True}
            raise FileNotFoundError(f"Path does not exist: {target}")
        if target.is_dir():
            if payload.recursive:
                shutil.rmtree(target)
            else:
                target.rmdir()
        else:
            target.unlink()
        return {"path": str(target), "deleted": True}

    def delete_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = DeleteManyInput.model_validate(arguments)
        return self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.delete_path(
                {
                    "path": raw_path,
                    "recursive": payload.recursive,
                    "missing_ok": payload.missing_ok,
                }
            ),
        )

    def open_path(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = OpenPathInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        if not target.exists():
            raise FileNotFoundError(f"Path does not exist: {target}")
        self._system_open(target)
        return {"path": str(target), "opened": True, "method": "default"}

    def open_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = PathListInput.model_validate(arguments)
        return self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.open_path({"path": raw_path}),
        )

    def reveal_in_explorer(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = OpenPathInput.model_validate(arguments)
        target = self.resolve_path(payload.path)
        self.ensure_workspace_path(target)
        if not target.exists():
            raise FileNotFoundError(f"Path does not exist: {target}")
        if os.name == "nt":
            if target.is_dir():
                subprocess.Popen(["explorer.exe", str(target)])
            else:
                subprocess.Popen(["explorer.exe", f"/select,{target}"])
        else:
            self._system_open(target.parent if target.is_file() else target)
        return {"path": str(target), "opened": True, "method": "reveal"}

    def reveal_many(self, arguments: dict[str, Any]) -> dict[str, Any]:
        payload = PathListInput.model_validate(arguments)
        return self._run_many_operations(
            payload.paths,
            continue_on_error=payload.continue_on_error,
            key_builder=lambda raw_path: raw_path,
            operation=lambda raw_path: self.reveal_in_explorer({"path": raw_path}),
        )

    def _metadata_for_path(self, target: Path) -> dict[str, Any]:
        stat = target.stat()
        mime_type, _ = guess_type(str(target))
        metadata = {
            "path": str(target),
            "name": target.name,
            "is_dir": target.is_dir(),
            "is_file": target.is_file(),
            "size_bytes": stat.st_size,
            "modified_at": stat.st_mtime,
            "suffix": target.suffix,
            "mime_type": mime_type,
        }
        if target.is_dir():
            try:
                metadata["child_count"] = sum(1 for _ in target.iterdir())
            except OSError:
                metadata["child_count"] = None
        return metadata

    @staticmethod
    def _tokenize_search_query(query: str) -> list[str]:
        lowered = query.lower()
        for token in ("-", "_", "/", "\\", ".", "(", ")", "[", "]"):
            lowered = lowered.replace(token, " ")
        return [term for term in lowered.split() if term]

    def _score_name_candidate(
        self,
        item: Path,
        *,
        query_text: str,
        query_terms: list[str],
        alias_terms: list[str],
        file_type_hints: list[str],
    ) -> tuple[float, str]:
        name = item.name.lower()
        stem = item.stem.lower()
        path_text = item.as_posix().lower()
        compact_name = "".join(ch for ch in stem if ch.isalnum() or ord(ch) > 127)
        file_type_aliases: set[str] = set()
        for ext in file_type_hints:
            aliases = FileQueryNormalizer._FILE_TYPE_ALIASES.get(ext, ())
            file_type_aliases.update(alias.lower() for alias in aliases)
        semantic_query_terms = [term for term in query_terms if term.lower() not in file_type_aliases]
        semantic_alias_terms = [term for term in alias_terms if term.lower() not in file_type_aliases]
        compact_query = "".join(semantic_query_terms) if semantic_query_terms else query_text.replace(" ", "")

        score = 0.0
        reason = "term_overlap"
        if query_text and query_text == stem:
            score = 1.0
            reason = "exact_stem_match"
        elif query_text and query_text in name:
            score = 0.96
            reason = "name_contains_query"
        elif compact_query and compact_query in compact_name:
            score = 0.92
            reason = "compact_name_match"
        else:
            term_hits = sum(1 for term in semantic_query_terms if term in name or term in path_text)
            alias_hits = sum(1 for term in semantic_alias_terms if term.lower() in name or term.lower() in path_text)
            if term_hits:
                score = 0.45 + 0.45 * (term_hits / max(len(semantic_query_terms), 1))
                reason = "term_overlap"
                if len(semantic_query_terms) >= 2 and term_hits >= max(2, len(semantic_query_terms) - 1):
                    score += 0.12
                    reason = "unordered_core_terms"
                if alias_hits:
                    score += 0.18 * (alias_hits / max(len(semantic_alias_terms), 1))
                    reason = "term_and_alias_overlap"
            elif alias_hits:
                score = 0.34 + 0.26 * (alias_hits / max(len(semantic_alias_terms), 1))
                reason = "alias_overlap"
            if compact_query and compact_name:
                fuzzy_ratio = SequenceMatcher(None, compact_query, compact_name).ratio()
                if fuzzy_ratio >= 0.68:
                    fuzzy_score = 0.32 + 0.46 * fuzzy_ratio
                    if fuzzy_score > score:
                        score = fuzzy_score
                        reason = "fuzzy_name_similarity"
            if score > 0 and file_type_hints and item.is_file() and item.suffix.lower() in file_type_hints:
                score += 0.08
            if item.is_dir():
                score += 0.03
            elif score > 0 and item.suffix.lower() in {".docx", ".pptx", ".xlsx", ".pdf", ".md", ".txt"}:
                score += 0.02
            if self._looks_like_generic_document_name(stem):
                score *= 0.45
                reason = "generic_document_name"
        return min(score, 1.0), reason

    @staticmethod
    def _looks_like_generic_document_name(stem: str) -> bool:
        generic_markers = {
            "new pptx presentation",
            "presentation",
            "untitled",
            "新建 pptx 演示文稿",
            "新建演示文稿",
            "课件",
        }
        return stem.strip().lower() in generic_markers

    def _extract_text_from_file(
        self,
        path: Path,
        encoding: str,
        max_chars: int,
        max_rows_per_sheet: int,
    ) -> tuple[str, str]:
        suffix = path.suffix.lower()
        if suffix == ".docx":
            return self._extract_docx_text(path, max_chars), "docx"
        if suffix == ".pptx":
            return self._extract_pptx_text(path, max_chars), "pptx"
        if suffix == ".xlsx":
            return self._extract_xlsx_text(path, max_chars, max_rows_per_sheet), "xlsx"
        return path.read_text(encoding=encoding)[:max_chars], "text"

    def _extract_structure_from_file(
        self,
        path: Path,
        *,
        encoding: str,
        include_text: bool,
        max_blocks: int,
        max_chars_per_block: int,
        max_rows_per_sheet: int,
    ) -> dict[str, Any]:
        suffix = path.suffix.lower()
        if suffix == ".docx":
            return self._extract_docx_structure(path, include_text=include_text, max_blocks=max_blocks, max_chars_per_block=max_chars_per_block)
        content, extraction_type = self._extract_text_from_file(
            path,
            encoding=encoding,
            max_chars=max_chars_per_block,
            max_rows_per_sheet=max_rows_per_sheet,
        )
        text = content if include_text else ""
        return {
            "file_type": extraction_type,
            "blocks": [
                {
                    "block_id": "text-1",
                    "block_type": "text",
                    "text": text,
                    "style_name": None,
                    "heading_level": None,
                }
            ],
            "block_count": 1,
            "template_profile": {"file_type": extraction_type, "block_count": 1},
        }

    def _extract_docx_structure(
        self,
        path: Path,
        *,
        include_text: bool,
        max_blocks: int,
        max_chars_per_block: int,
    ) -> dict[str, Any]:
        from docx import Document

        document = Document(path)
        blocks: list[dict[str, Any]] = []
        style_counts: dict[str, int] = {}
        table_count = 0
        for block_ref in self._iter_docx_block_refs(document):
            if len(blocks) >= max_blocks:
                break
            text = str(block_ref.get("text", "")).strip()
            if not text:
                continue
            if block_ref.get("block_type") == "table_cell":
                table_count = max(table_count, int(block_ref.get("table_index", 0)) + 1)
            style_name = block_ref.get("style_name")
            if isinstance(style_name, str) and style_name:
                style_counts[style_name] = style_counts.get(style_name, 0) + 1
            blocks.append(
                {
                    "block_id": block_ref["block_id"],
                    "block_type": block_ref["block_type"],
                    "text": text[:max_chars_per_block] if include_text else "",
                    "style_name": style_name,
                    "heading_level": block_ref.get("heading_level"),
                    "table_index": block_ref.get("table_index"),
                    "row_index": block_ref.get("row_index"),
                    "col_index": block_ref.get("col_index"),
                }
            )
        return {
            "file_type": "docx",
            "blocks": blocks,
            "block_count": len(blocks),
            "template_profile": {
                "file_type": "docx",
                "style_counts": style_counts,
                "heading_count": sum(1 for block in blocks if block.get("heading_level")),
                "table_count": table_count,
                "block_count": len(blocks),
            },
        }

    @staticmethod
    def _extract_docx_text(path: Path, max_chars: int) -> str:
        with ZipFile(path) as archive:
            xml_text = archive.read("word/document.xml").decode("utf-8", errors="ignore")
        root = ET.fromstring(xml_text)
        namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        paragraphs = []
        for paragraph in root.findall(".//w:p", namespace):
            texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
            line = "".join(texts).strip()
            if line:
                paragraphs.append(line)
        return "\n".join(paragraphs)[:max_chars]

    @staticmethod
    def _extract_pptx_text(path: Path, max_chars: int) -> str:
        with ZipFile(path) as archive:
            slide_names = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
            slides: list[str] = []
            namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for index, slide_name in enumerate(slide_names, start=1):
                xml_text = archive.read(slide_name).decode("utf-8", errors="ignore")
                root = ET.fromstring(xml_text)
                texts = [node.text or "" for node in root.findall(".//a:t", namespace)]
                line = " ".join(part.strip() for part in texts if part and part.strip()).strip()
                if line:
                    slides.append(f"Slide {index}: {line}")
            return "\n".join(slides)[:max_chars]

    @staticmethod
    def _extract_xlsx_text(path: Path, max_chars: int, max_rows_per_sheet: int) -> str:
        namespace_main = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        namespace_rel = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
        with ZipFile(path) as archive:
            shared_strings: list[str] = []
            if "xl/sharedStrings.xml" in archive.namelist():
                shared_root = ET.fromstring(archive.read("xl/sharedStrings.xml").decode("utf-8", errors="ignore"))
                shared_strings = [
                    "".join(node.itertext()).strip()
                    for node in shared_root.findall(".//main:si", namespace_main)
                ]

            workbook_root = ET.fromstring(archive.read("xl/workbook.xml").decode("utf-8", errors="ignore"))
            rel_root = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels").decode("utf-8", errors="ignore"))
            rel_map = {
                rel.attrib.get("Id"): rel.attrib.get("Target", "")
                for rel in rel_root.findall(".//rel:Relationship", namespace_rel)
            }

            lines: list[str] = []
            for sheet in workbook_root.findall(".//main:sheets/main:sheet", namespace_main):
                sheet_name = sheet.attrib.get("name", "Sheet")
                rel_id = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                target = rel_map.get(rel_id, "")
                if not target:
                    continue
                sheet_path = "xl/" + target.lstrip("/")
                sheet_root = ET.fromstring(archive.read(sheet_path).decode("utf-8", errors="ignore"))
                lines.append(f"[Sheet] {sheet_name}")
                for row in sheet_root.findall(".//main:sheetData/main:row", namespace_main)[:max_rows_per_sheet]:
                    values: list[str] = []
                    for cell in row.findall("main:c", namespace_main):
                        cell_type = cell.attrib.get("t")
                        value_node = cell.find("main:v", namespace_main)
                        inline_node = cell.find("main:is/main:t", namespace_main)
                        value = ""
                        if cell_type == "s" and value_node is not None and value_node.text and value_node.text.isdigit():
                            idx = int(value_node.text)
                            if 0 <= idx < len(shared_strings):
                                value = shared_strings[idx]
                        elif inline_node is not None and inline_node.text:
                            value = inline_node.text
                        elif value_node is not None and value_node.text:
                            value = value_node.text
                        if value:
                            values.append(value)
                    if values:
                        lines.append(" | ".join(values))
            return "\n".join(lines)[:max_chars]

    @staticmethod
    def _split_paragraphs(content: str) -> list[str]:
        return [paragraph.strip() for paragraph in content.splitlines() if paragraph.strip()]

    @classmethod
    def _build_tabular_rows(cls, content: str) -> list[list[str]]:
        rows: list[list[str]] = []
        for line in cls._split_paragraphs(content):
            if "," in line:
                cells = [cell.strip() for cell in line.split(",")]
            elif "\t" in line:
                cells = [cell.strip() for cell in line.split("\t")]
            elif ":" in line:
                key, value = line.split(":", maxsplit=1)
                cells = [key.strip(), value.strip()]
            else:
                cells = [line]
            if any(cell for cell in cells):
                rows.append(cells)
        return rows

    @classmethod
    def _build_bullets(cls, content: str) -> list[str]:
        bullets: list[str] = []
        for line in cls._split_paragraphs(content):
            cleaned = line.lstrip("-*0123456789. ").strip()
            if cleaned:
                bullets.append(cleaned)
        return bullets

    @staticmethod
    def _docx_heading_level(style_name: str | None) -> int | None:
        if not style_name:
            return None
        match = re.search(r"heading\s*(\d+)", style_name, flags=re.IGNORECASE)
        if match:
            return int(match.group(1))
        return None

    @classmethod
    def _iter_docx_block_refs(cls, document) -> list[dict[str, Any]]:
        from docx.document import Document as DocxDocument
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph

        def iter_block_items(parent):
            if isinstance(parent, DocxDocument):
                parent_elm = parent.element.body
            elif isinstance(parent, _Cell):
                parent_elm = parent._tc
            else:
                raise TypeError("Unsupported parent type for docx traversal")
            for child in parent_elm.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)

        blocks: list[dict[str, Any]] = []
        paragraph_index = 0
        table_index = 0
        for item in iter_block_items(document):
            if isinstance(item, Paragraph):
                paragraph_index += 1
                text = item.text.strip()
                if not text:
                    continue
                style_name = item.style.name if item.style is not None else None
                blocks.append(
                    {
                        "block_id": f"p-{paragraph_index}",
                        "block_type": "paragraph",
                        "text": text,
                        "style_name": style_name,
                        "heading_level": cls._docx_heading_level(style_name),
                        "paragraph": item,
                    }
                )
            elif isinstance(item, Table):
                current_table_index = table_index
                table_index += 1
                for row_index, row in enumerate(item.rows):
                    for col_index, cell in enumerate(row.cells):
                        cell_text = "\n".join(part.strip() for part in cell.text.splitlines() if part.strip()).strip()
                        if not cell_text:
                            continue
                        first_paragraph = cell.paragraphs[0] if cell.paragraphs else None
                        style_name = None
                        if first_paragraph is not None and first_paragraph.style is not None:
                            style_name = first_paragraph.style.name
                        blocks.append(
                            {
                                "block_id": f"t-{current_table_index}-r-{row_index}-c-{col_index}",
                                "block_type": "table_cell",
                                "text": cell_text,
                                "style_name": style_name,
                                "heading_level": cls._docx_heading_level(style_name),
                                "table_index": current_table_index,
                                "row_index": row_index,
                                "col_index": col_index,
                                "cell": cell,
                            }
                        )
        return blocks

    @classmethod
    def _build_docx_block_reference_map(cls, document) -> dict[str, dict[str, Any]]:
        return {block["block_id"]: block for block in cls._iter_docx_block_refs(document)}

    @staticmethod
    def _replace_docx_block_text(block_ref: dict[str, Any], text: str) -> None:
        if block_ref.get("block_type") == "table_cell":
            cell = block_ref["cell"]
            cell.text = text
            return
        paragraph = block_ref["paragraph"]
        paragraph.text = text

    @staticmethod
    def _delete_docx_block(block_ref: dict[str, Any]) -> None:
        if block_ref.get("block_type") == "table_cell":
            block_ref["cell"].text = ""
            return
        paragraph = block_ref["paragraph"]
        element = paragraph._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    @staticmethod
    def _insert_docx_block_after(block_ref: dict[str, Any], text: str) -> None:
        if block_ref.get("block_type") == "table_cell":
            cell = block_ref["cell"]
            new_paragraph = cell.add_paragraph(text)
            style_name = block_ref.get("style_name")
            if style_name:
                try:
                    new_paragraph.style = style_name
                except Exception:  # noqa: BLE001
                    pass
            return
        from docx.oxml import OxmlElement
        from docx.text.paragraph import Paragraph

        paragraph = block_ref["paragraph"]
        new_element = OxmlElement("w:p")
        paragraph._p.addnext(new_element)
        new_paragraph = Paragraph(new_element, paragraph._parent)
        style_name = block_ref.get("style_name")
        if style_name:
            try:
                new_paragraph.style = style_name
            except Exception:  # noqa: BLE001
                pass
        new_paragraph.text = text

    def _resolve_render_docx_paragraphs(
        self,
        payload: RenderDocxFromTemplateInput,
        source_path: Path | None,
    ) -> list[str]:
        if payload.paragraphs:
            return [paragraph.strip() for paragraph in payload.paragraphs if paragraph.strip()]
        if payload.content.strip():
            return self._split_paragraphs(payload.content)
        if source_path is None:
            return []
        structure = self._extract_structure_from_file(
            source_path,
            encoding="utf-8",
            include_text=True,
            max_blocks=400,
            max_chars_per_block=1200,
            max_rows_per_sheet=10,
        )
        return [
            str(block.get("text", "")).strip()
            for block in structure.get("blocks", [])
            if isinstance(block, dict) and str(block.get("text", "")).strip()
        ]

    def _remove_existing_path(self, target: Path) -> None:
        if not target.exists():
            return
        if target.is_dir():
            shutil.rmtree(target)
        else:
            target.unlink()

    @staticmethod
    def _run_many_operations(
        items: list[Any],
        *,
        continue_on_error: bool,
        key_builder,
        operation,
    ) -> dict[str, Any]:
        results: list[dict[str, Any]] = []
        successful_paths: list[str] = []
        failure_count = 0
        for item in items:
            key = key_builder(item)
            try:
                result = operation(item)
                entry = dict(result)
                entry["ok"] = True
                results.append(entry)
                path = entry.get("path")
                if isinstance(path, str) and path:
                    successful_paths.append(path)
            except Exception as exc:  # noqa: BLE001
                if not continue_on_error:
                    raise
                failure_count += 1
                results.append({"path": str(key), "ok": False, "error": str(exc)})
        return {
            "paths": successful_paths,
            "results": results,
            "success_count": len(successful_paths),
            "failure_count": failure_count,
        }

    @staticmethod
    def _system_open(target: Path) -> None:
        if hasattr(os, "startfile"):
            os.startfile(str(target))  # type: ignore[attr-defined]
            return
        subprocess.Popen(["xdg-open", str(target)])
